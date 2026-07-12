from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG=RGBColor(0x0a,0x12,0x11); INK=RGBColor(0xff,0xff,0xff); MUT=RGBColor(0x8f,0xa3,0x9d)
TEAL=RGBColor(0x2f,0xd6,0xbf); AMBER=RGBColor(0xf4,0xb0,0x62)
W,H=Inches(13.333),Inches(7.5)

# eyebrow, title (parts: (text,color)), narration note
SLIDES=[
 ("BRAKEPOINT · BUILT WITH CLAUDE",[("The ",INK),("brakes",AMBER),(" on a T cell are its drug targets.",INK)],
  "Checkpoint immunotherapy works by releasing the brakes on a T cell. So we went looking for those brakes — genome-wide. This is Brakepoint, built with Claude Science."),
 ("THE THESIS",[("Which knockdowns make a T cell a ",INK),("stronger effector",AMBER),("?",INK)],
  "Those brakes matter beyond checkpoint therapy — they also throttle engineered CAR-T. So across the genome, which genes are the brakes — which knockdowns push a human T cell toward a stronger effector state?"),
 ("THE SCREEN",[("2.6M CD4",INK),("+",TEAL),(" T cells · 12,449 CRISPRi knockdowns.",INK)],
  "We started from a genome-scale CRISPR-interference screen: over twelve thousand gene knockdowns, across 2.6 million primary human CD4 T cells, from the Gladstone Institutes. Two donors, out of an intended four."),
 ("THE TRAP · SIGNIFICANCE",[("Significance can't rank a ",INK),("million-cell screen",TEAL),(".",INK)],
  "How do you find the brakes in twelve thousand knockdowns? The reflex is to rank by significance — but at two million cells that breaks down: over 97% of the tested knockdowns clear the bar. So we rank by causal effect size instead."),
 ("THE ENGINE · SIGNED CAUSAL MAP",[("Rank by effect size, then add the ",INK),("sign",AMBER),(".",INK)],
  "But effect size alone still points at the wrong genes — the biggest hits are the cell's own signaling machinery. So we add what a magnitude ranking leaves out: direction — toward the effector program, or away. Now the machinery falls away, and the candidate brakes rise into view."),
 ("WHY IT'S DIFFERENT",[("Same data. ",INK),("A sharper question.",AMBER)],
  "And that's the edge. Differential expression finds correlations — not what to drug. Genetics points to a locus, rarely a direction. We measure what a knockdown actually does, and which way it pushes — then weigh it against genetics and the clinic."),
 ("THE SHORTLIST",[("Five candidate brakes: ",INK),("CBLB, CD5, DGKA, SMAD3, UBASH3A",AMBER),(".",INK)],
  "From that map, five prior-informed candidates: CBLB, CD5, DGKA, SMAD3, and UBASH3A — each scored across seven lines of evidence, from causal effect and direction to human genetics and clinical precedent."),
 ("THE LEAD · CBLB",[("CBLB",AMBER),(" — inhibitors already in early-phase trials.",INK)],
  "Our lead is CBLB. It's a natural off-switch for T-cell activation, and inhibitors are already in early-phase trials. Its genetics point the same way, and it lands in our brake quadrant. CD5 and DGKA come next — and both hold up across donors."),
 ("REPORTED HONESTLY",[("CD5 & DGKA hold up; the rest are ",INK),("donor-split",AMBER),(" (p = 0.70).",INK)],
  "And we're honest about what two donors can support. CD5 and DGKA hold up in both; the other three ride on one donor. As a group, known brakes aren't significantly enriched — so this is a ranked shortlist for the full cohort, a hypothesis to test, not a finished answer."),
 ("OPEN SOURCE · MIT",[("Every number traces back to ",INK),("code",TEAL),(".",INK)],
  "Every number here traces back to versioned code — and to a self-check that caught a real bias in our effect-size code before it reached a figure. Open source; every figure regenerates with one command. This is Brakepoint, built with Claude Science."),
]

prs=Presentation(); prs.slide_width=W; prs.slide_height=H
blank=prs.slide_layouts[6]
for eyebrow,parts,note in SLIDES:
    s=prs.slides.add_slide(blank)
    # bg
    s.background.fill.solid(); s.background.fill.fore_color.rgb=BG
    # eyebrow
    eb=s.shapes.add_textbox(Inches(0.9),Inches(1.0),Inches(11.5),Inches(0.5)).text_frame
    r=eb.paragraphs[0].add_run(); r.text=eyebrow; f=r.font; f.size=Pt(15); f.bold=True; f.color.rgb=TEAL; f.name="Arial"
    # title
    tb=s.shapes.add_textbox(Inches(0.9),Inches(2.2),Inches(11.5),Inches(3.2)).text_frame; tb.word_wrap=True
    p=tb.paragraphs[0]
    for txt,col in parts:
        r=p.add_run(); r.text=txt; f=r.font; f.size=Pt(46); f.bold=True; f.color.rgb=col; f.name="Arial"
    # footer wordmark
    fb=s.shapes.add_textbox(Inches(0.9),Inches(6.7),Inches(11.5),Inches(0.5)).text_frame
    r=fb.paragraphs[0].add_run(); r.text="Brakepoint · Built with Claude Science"; f=r.font; f.size=Pt(12); f.color.rgb=MUT; f.name="Arial"
    # speaker notes = narration
    s.notes_slide.notes_text_frame.text=note
prs.save("demo_deck.pptx")
print("rebuilt demo_deck.pptx:", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
